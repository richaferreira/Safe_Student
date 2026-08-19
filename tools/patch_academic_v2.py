from pathlib import Path
from docx import Document
from docx.shared import Cm, Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.section import WD_ORIENT
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.enum.shapes import MSO_SHAPE_TYPE


def replace_in_docx(path, replacements):
    doc = Document(path)

    def replace_para(p):
        full = ''.join(r.text for r in p.runs)
        new = full
        for old, repl in replacements:
            new = new.replace(old, repl)
        if new != full:
            p.text = new
            for r in p.runs:
                r.font.name = 'Times New Roman'
                r.font.size = Pt(12)
            return 1
        return 0

    changed = 0
    for p in doc.paragraphs:
        changed += replace_para(p)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    changed += replace_para(p)
    doc.save(path)
    print(f'{path}: {changed} substituições em parágrafos/células')


def patch_plan():
    path = Path('01_Plano_Gerenciamento_Unificado_Safe_Student.docx')
    doc = Document(path)
    # Corrige a tabela da pesquisa sem fabricar nova evidência.
    replacements_by_label = {
        'Participantes da validação': 'Resultados históricos declarados na versão anterior; a evidência primária não foi localizada no repositório. Não apresentar quantitativo como validado sem anexar formulários ou CSV verificável.',
        'Tarefa 1 - localizar histórico': 'Pendente de comprovação primária ou de nova coleta controlada.',
        'Tarefa 2 - registrar presença': 'Pendente de comprovação primária ou de nova coleta controlada.',
        'Tarefa 3 - compreender notificação': 'Pendente de comprovação primária ou de nova coleta controlada.',
        'Satisfação geral': 'Pendente de comprovação primária. A nova coleta deve registrar nota individual no módulo Feedback.',
        'Principal melhoria sugerida': 'A nova coleta deve registrar cenário, conclusão da tarefa, tempo, nota e comentário; somente dados APRESENTACAO podem ser tratados como evidência coletada.',
    }
    for table in doc.tables:
        for row in table.rows:
            if len(row.cells) < 2:
                continue
            label = row.cells[0].text.strip()
            if label in replacements_by_label:
                row.cells[1].text = replacements_by_label[label]
                for p in row.cells[1].paragraphs:
                    for r in p.runs:
                        r.font.name = 'Times New Roman'; r.font.size = Pt(9)
    for p in doc.paragraphs:
        text = ''.join(r.text for r in p.runs)
        new = text.replace('Validar 5 funcionalidades essenciais e 15 RF/RNF com testes automatizados e avaliação acadêmica controlada de usabilidade.', 'Validar as 5 funcionalidades essenciais por requisitos rastreáveis, testes automatizados e avaliação acadêmica controlada de usabilidade.')
        new = new.replace('15 RF/RNF', 'requisitos funcionais e não funcionais rastreáveis')
        if new != text:
            p.text = new
            for r in p.runs:
                r.font.name = 'Times New Roman'; r.font.size = Pt(12)
    doc.save(path)
    print('Plano V2: evidência de pesquisa e contagem artificial de requisitos corrigidas.')


def patch_test_report():
    path = Path('04_Plano_Relatorio_Testes_Safe_Student.docx')
    doc = Document(path)
    for p in doc.paragraphs:
        text = ''.join(r.text for r in p.runs)
        if '19' in text and any(k in text.lower() for k in ['teste', 'caso', 'suíte', 'suite']):
            text = text.replace('19 testes automatizados', 'suíte automatizada atual')
            text = text.replace('19 casos automatizados', 'suíte automatizada atual')
            text = text.replace('19 casos', 'casos automatizados atuais')
            text = text.replace('19 testes', 'testes automatizados atuais')
            p.text = text
            for r in p.runs:
                r.font.name = 'Times New Roman'; r.font.size = Pt(12)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text = cell.text
                if '19' in text and any(k in text.lower() for k in ['teste','caso','suíte','suite']):
                    cell.text = text.replace('19 testes automatizados','suíte automatizada atual').replace('19 casos automatizados','suíte automatizada atual').replace('19 casos','casos automatizados atuais').replace('19 testes','testes automatizados atuais')
    doc.save(path)
    print('Relatório de testes: removida dependência de contagem fixa que fica obsoleta com novos testes.')


def build_audit_report():
    path = Path('06_Relatorio_Auditoria_Conformidade_Safe_Student.docx')
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Cm(3); sec.left_margin = Cm(3); sec.right_margin = Cm(2); sec.bottom_margin = Cm(2)
    doc.styles['Normal'].font.name = 'Times New Roman'; doc.styles['Normal'].font.size = Pt(12)
    for level in [1,2]:
        st=doc.styles[f'Heading {level}']; st.font.name='Times New Roman'; st.font.bold=True; st.font.size=Pt(14 if level==1 else 12)
    for _ in range(3): doc.add_paragraph()
    p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER; r=p.add_run('SAFE STUDENT'); r.bold=True; r.font.size=Pt(18)
    p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER; r=p.add_run('RELATÓRIO DE AUDITORIA TÉCNICA E DE CONFORMIDADE - V2.0'); r.bold=True; r.font.size=Pt(15)
    p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER; p.add_run('Revisão com foco em requisitos, UML/DER, coerência com o código, evidência e qualidade de engenharia')
    doc.add_page_break()

    def h(text,level=1): doc.add_paragraph(text,style=f'Heading {level}')
    def body(text):
        p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.JUSTIFY; p.paragraph_format.first_line_indent=Cm(1.25); p.paragraph_format.line_spacing=1.5; p.add_run(text)
    def table(headers,rows):
        t=doc.add_table(rows=1,cols=len(headers)); t.style='Table Grid'; t.alignment=WD_TABLE_ALIGNMENT.CENTER
        for i,v in enumerate(headers): t.rows[0].cells[i].text=v
        for row in rows:
            c=t.add_row().cells
            for i,v in enumerate(row): c[i].text=str(v)
        for ri,row in enumerate(t.rows):
            for cell in row.cells:
                for p in cell.paragraphs:
                    for r in p.runs: r.font.name='Times New Roman'; r.font.size=Pt(9); r.bold=(ri==0)
        return t

    h('1 OBJETIVO DA AUDITORIA')
    body('A auditoria V2.0 foi feita com postura de análise de sistemas e desenvolvimento sênior: cada afirmação documental foi confrontada com o comportamento observável do MVP, com especial atenção a requisitos testáveis, autorização no servidor, privacidade, persistência, rastreabilidade e uso correto da notação UML. O relatório não transforma metas futuras em funcionalidades implementadas e não considera números de pesquisa como evidência quando os registros primários não estão disponíveis.')
    h('2 ACHADOS CRÍTICOS E CORREÇÕES')
    rows=[
        ('Crítico','Diagramas de casos de uso com excesso de cruzamentos e linhas que pareciam setas de fluxo.','Diagramas separados por ator; associações ator-caso de uso usam linha sólida simples sem seta.'),
        ('Crítico','Modelo de classes não distinguia claramente conceito, implementação e relacionamentos.','Classes reorganizadas; Usuario permanece único com atributo perfil; multiplicidades explícitas.'),
        ('Crítico','DER visualmente confundido com fluxograma.','DER lógico refeito com PK/FK/UQ e cardinalidades; relacionamentos não usam setas.'),
        ('Alto','Documento citava OWASP MASVS, que é orientado a aplicações móveis, para um MVP web.','Referência trocada por OWASP ASVS 5.0.0.'),
        ('Alto','Gestão podia receber notificações e visualizar mensagens de terceiros por regra ampla no backend.','Consultas passam a respeitar destinatário/participante; testes de regressão adicionados.'),
        ('Alto','Diretório retornava mais dados de usuário do que a interface precisava.','Projeção mínima: id, nome, perfil e status.'),
        ('Alto','Exportação de validação incluía registros DEMO_SEED.','CSV passa a exportar somente registros APRESENTACAO.'),
        ('Médio','Taxa de frequência usava a data UTC do timestamp, diferente do conceito de dia escolar do restante do servidor.','attendanceRate passa a aceitar a função dateKey com fuso escolar.'),
        ('Médio','Token de novo aluno era aleatório, mas não havia verificação explícita de colisão.','Geração agora verifica unicidade antes de aceitar o token.'),
        ('Médio','db.json mutável ficou versionado e recebeu estado de execução.','Execução usa db.runtime.json, criado a partir da seed e ignorado pelo Git.'),
        ('Médio','Workflows antigos ficaram quebrados após a exclusão dos scripts que eles chamavam.','Workflows obsoletos removidos; CI principal permanece.'),
        ('Acadêmico','Plano/apresentação mostravam 15 participantes e média 4,4 sem evidência primária disponível.','Números deixam de ser apresentados como resultados auditados; nova coleta deve gerar evidência real.')]
    table(['Severidade','Achado','Correção'],rows)
    h('3 REQUISITOS')
    body('Os requisitos foram reescritos para descrever obrigações do sistema, não telas ou intenções vagas. Cada RF possui critério de aceite. Os RNF indicam forma de verificação e não declaram cumprimento quando não existe medição. O painel de privacidade foi tratado como recurso de apoio da interface, e não como requisito funcional central de negócio.')
    h('4 MODELAGEM UML E DER')
    body('A especificação V2.0 usa a notação UML de modo conservador. Em casos de uso, ator e caso de uso são ligados por associação sem seta. No diagrama de classes, associações usam multiplicidade e não são transformadas em setas de processo. No diagrama de sequência, setas são apropriadas porque representam mensagens dirigidas. No DER, as ligações representam relacionamentos e são acompanhadas por chaves e cardinalidades.')
    h('5 AUDITORIA DO CÓDIGO')
    table(['Arquivo/área','Conclusão'],[
        ('server.js','Revisado em autenticação, autorização, escopo de mensagens/notificações, minimização de diretório, token único, reset auditado e exportação de evidências.'),
        ('lib/domain.js','Regras de perfil e sequência coerentes; cálculo de taxa ajustado para receber o dia escolar.'),
        ('lib/security.js','Adequado ao MVP: scrypt + timingSafeEqual. scryptSync é aceitável para demonstração, mas não é recomendação de arquitetura de alta carga.'),
        ('public/app.js','Consome API autenticada e escapa conteúdo renderizado. Controles visuais não são tratados como barreira de segurança; a autorização permanece no backend.'),
        ('public/index.html','Estrutura possui rótulos, skip-link e aviso de ambiente de demonstração. Não deve ser apresentado como certificação WCAG.'),
        ('data/db.seed.json','Dados sintéticos de demonstração. Não equivalem a evidência de pesquisa de campo.'),
        ('tests','Além dos testes existentes, foram adicionadas regressões de privacidade, isolamento de mensagens/notificações, exportação de feedback e auditoria do reset.')])
    h('6 DOCUMENTAÇÃO E EVIDÊNCIAS')
    body('A pendência acadêmica que continua dependendo de ação humana é a comprovação primária da pesquisa/validação com participantes. Sem formulários anonimizados, termos preenchidos ou CSV de coleta real, os resultados históricos não devem ser tratados como auditados. A versão revisada mantém a distinção entre dados DEMO_SEED e dados APRESENTACAO.')
    h('7 SITUAÇÃO APÓS A REVISÃO')
    table(['Dimensão','Status'],[
        ('Código do MVP','Revisado; aguardando/condicionado a CI verde na branch de revisão antes de integrar à main.'),
        ('Requisitos','Reescritos e rastreáveis.'),
        ('UML/DER','Refeito com notação formal e figuras separadas para legibilidade.'),
        ('Documentação acadêmica','Corrigida onde havia afirmações não comprovadas ou referências técnicas inadequadas.'),
        ('Pesquisa de campo','Parcial até anexar evidência primária real.'),
        ('Produção','Fora do escopo: o MVP continua sendo demonstração local e não solução homologada.')])
    h('REFERÊNCIAS')
    for ref in ['OBJECT MANAGEMENT GROUP. Unified Modeling Language (UML), Version 2.5.1.','OWASP FOUNDATION. Application Security Verification Standard (ASVS) 5.0.0.','WORLD WIDE WEB CONSORTIUM. Web Content Accessibility Guidelines (WCAG) 2.2.','BRASIL. Lei nº 13.709/2018 - Lei Geral de Proteção de Dados Pessoais.']:
        doc.add_paragraph(ref)
    doc.save(path)
    print('Relatório de auditoria V2 reconstruído.')


def patch_manual():
    replace_in_docx(Path('10_Manual_Usuario_Implantacao_Treinamento_Suporte.docx'), [('db.json','db.runtime.json'),('OWASP MASVS','OWASP ASVS')])


def patch_slides():
    path=Path('08_Apresentacao_Final_Safe_Student.pptx')
    prs=Presentation(path)
    replacements={
        '15 participantes':'Evidência primária pendente',
        'Nota média 4,4':'Nova coleta estruturada',
        '3 perfis':'3 perfis previstos',
        '30 requisitos estruturam o sistema e conectam documentação, MVP e testes.':'Requisitos revisados conectam comportamento, código e testes.',
        '15 RF':'RF verificáveis',
        '15 RNF':'RNF verificáveis',
        'OWASP MASVS':'OWASP ASVS 5.0.0',
        'pacote acadêmico-profissional':'pacote acadêmico revisado e demonstrável',
    }
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape,'has_text_frame',False):
                continue
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    for old,new in replacements.items():
                        if old in run.text: run.text=run.text.replace(old,new)
    if len(prs.slides) >= 6:
        slide=prs.slides[5]
        for shape in list(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                sp=shape._element; sp.getparent().remove(sp)
        y=PInches(2.0); width=PInches(4.0)
        imgs=['docs/diagramas/02_uc_portaria.png','docs/diagramas/04_classes_nucleo.png','docs/diagramas/06_der_nucleo.png']
        xs=[PInches(0.3),PInches(4.65),PInches(9.0)]
        caps=['Casos de uso','Classes','DER lógico']
        for x,img,cap in zip(xs,imgs,caps):
            slide.shapes.add_picture(img,x,y,width=width)
            tb=slide.shapes.add_textbox(x,PInches(6.6),width,PInches(0.35))
            p=tb.text_frame.paragraphs[0]; p.text=cap; p.alignment=1
            for r in p.runs: r.font.size=PPt(12); r.font.bold=True
    prs.save(path)
    print('Apresentação V2 corrigida: evidências, ASVS e modelagem.')


patch_plan()
patch_test_report()
build_audit_report()
patch_manual()
patch_slides()
