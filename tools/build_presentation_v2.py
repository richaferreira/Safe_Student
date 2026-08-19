from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = Path('.')
DIAG = Path('docs/diagramas')
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
NAVY = RGBColor(20, 42, 74)
TEAL = RGBColor(24, 137, 153)
DARK = RGBColor(31, 41, 55)
MID = RGBColor(87, 99, 116)
LIGHT = RGBColor(244, 247, 250)
BORDER = RGBColor(214, 222, 230)
WHITE = RGBColor(255, 255, 255)
WARN = RGBColor(156, 95, 24)


def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill=WHITE, line=BORDER, radius=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line; shp.line.width = Pt(1)
    return shp


def txt(slide, text, x, y, w, h, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font='Aptos', valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return box


def title(slide, kicker, headline, subtitle=None):
    txt(slide, kicker.upper(), 0.6, 0.35, 5.6, 0.3, 11, True, TEAL)
    txt(slide, headline, 0.6, 0.72, 12.0, 0.65, 26, True, NAVY)
    if subtitle:
        txt(slide, subtitle, 0.6, 1.42, 12.0, 0.45, 13, False, MID)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.95), Inches(12.1), Inches(0.02))
    bar.fill.solid(); bar.fill.fore_color.rgb = BORDER; bar.line.fill.background()


def footer(slide, n):
    txt(slide, 'Safe Student • Práticas Extensionistas Integradoras VI • Prof. Altemar Sales de Oliveira', 0.6, 7.12, 11.7, 0.2, 8, False, MID)
    txt(slide, str(n), 12.25, 7.08, 0.45, 0.2, 8, False, MID, PP_ALIGN.RIGHT)


def bullet_card(slide, x, y, w, h, head, body, accent=TEAL):
    rect(slide, x, y, w, h, WHITE, BORDER, True)
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = accent; shp.line.fill.background()
    txt(slide, head, x + 0.22, y + 0.2, w - 0.4, 0.32, 15, True, NAVY)
    txt(slide, body, x + 0.22, y + 0.63, w - 0.4, h - 0.78, 11.5, False, DARK)


def metric(slide, x, y, w, value, label, accent=TEAL):
    rect(slide, x, y, w, 1.0, LIGHT, BORDER, True)
    txt(slide, value, x + 0.15, y + 0.13, w - 0.3, 0.4, 21, True, accent)
    txt(slide, label, x + 0.15, y + 0.58, w - 0.3, 0.25, 10.5, False, MID)


# 1
s = blank(); rect(s, 0, 0, 13.333, 7.5, WHITE, WHITE)
sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), Inches(7.5)); sh.fill.solid(); sh.fill.fore_color.rgb = TEAL; sh.line.fill.background()
txt(s, 'SAFE STUDENT', 0.85, 1.1, 11.6, 0.6, 34, True, NAVY)
txt(s, 'Presença escolar rastreável e comunicação família–escola', 0.85, 1.83, 11.3, 0.5, 20, False, DARK)
txt(s, 'Revisão técnica V2.0 • requisitos verificáveis • UML/DER formal • MVP auditável', 0.85, 2.52, 11.3, 0.45, 14, True, TEAL)
rect(s, 0.85, 3.4, 11.3, 1.55, LIGHT, BORDER, True)
txt(s, 'Flavio Gabrig Ferreira — 202323361\nRichardson Conceição Ferreira — 202323181\nGerente do projeto: Richardson Conceição Ferreira', 1.15, 3.72, 7.6, 0.95, 13, False, DARK)
txt(s, 'Práticas Extensionistas Integradoras VI\nProfessor: Altemar Sales de Oliveira\nSaquarema/RJ • 2026', 9.0, 3.72, 2.8, 0.95, 12, False, DARK, PP_ALIGN.RIGHT)
txt(s, 'MVP acadêmico com dados sintéticos — não homologado para produção', 0.85, 6.6, 11.3, 0.35, 11, False, MID)

# 2
s = blank(); title(s, 'Contexto', 'Problema, proposta de valor e limite do MVP')
bullet_card(s, 0.65, 2.3, 3.8, 2.0, 'Problema', 'Processos manuais e comunicação dispersa dificultam confirmar entrada/saída, consolidar histórico e responder rapidamente a dúvidas da família.')
bullet_card(s, 4.75, 2.3, 3.8, 2.0, 'Proposta', 'Identificação simples por token, registro rastreável, notificação interna e consulta por perfil — com o mínimo de dados necessário.')
bullet_card(s, 8.85, 2.3, 3.8, 2.0, 'Limite acadêmico', 'Sem biometria, NFC/RFID, catraca, geolocalização contínua, push real ou banco corporativo. Esses itens não são desenhados como implementados.', WARN)
metric(s, 0.65, 4.75, 3.8, '5', 'funcionalidades essenciais'); metric(s, 4.75, 4.75, 3.8, '3', 'atores centrais: família, portaria e gestão'); metric(s, 8.85, 4.75, 3.8, '1', 'fonte de verdade técnica: comportamento verificável'); footer(s, 2)

# 3
s = blank(); title(s, 'Escopo', 'As cinco funcionalidades essenciais do MVP', 'Recursos complementares são mantidos fora da contagem para evitar inflar o escopo.')
items = [('1','Autenticação e autorização','Acesso e operações controladas por perfil.'),('2','Vínculo aluno–responsável','Família enxerga apenas estudantes vinculados.'),('3','Entrada e saída','Registro por token/QR simulado com sequência válida.'),('4','Notificação','Aviso interno criado para responsáveis vinculados.'),('5','Rastreabilidade','Histórico, relatório e auditoria de operações relevantes.')]
for i, (n, h1, body) in enumerate(items):
    y = 2.3 + i * 0.88; txt(s, n, 0.8, y, 0.45, 0.42, 17, True, TEAL, PP_ALIGN.CENTER); rect(s, 1.45, y - 0.05, 10.9, 0.67, WHITE, BORDER, True); txt(s, h1, 1.7, y + 0.06, 3.1, 0.28, 13.5, True, NAVY); txt(s, body, 4.75, y + 0.06, 7.2, 0.28, 11.5, False, DARK)
footer(s, 3)

# 4
s = blank(); title(s, 'Engenharia de requisitos', 'Requisitos V2: verificáveis, rastreáveis e coerentes com o código')
bullet_card(s, 0.65, 2.25, 3.85, 1.65, 'RF verificáveis', '15 requisitos funcionais descrevem obrigações observáveis. Cada RF possui critério de aceite e prioridade.')
bullet_card(s, 4.75, 2.25, 3.85, 1.65, 'RNF verificáveis', '12 requisitos não funcionais indicam a forma de comprovação. Metas sem medição não são tratadas como atendidas.')
bullet_card(s, 8.85, 2.25, 3.85, 1.65, '12 regras de negócio', 'Sequência de presença, escopo por vínculo, perfis autorizados, evidência de feedback e definição do dia escolar.')
rect(s, 0.65, 4.35, 12.05, 1.55, LIGHT, BORDER, True); txt(s, 'Mudança de postura', 0.95, 4.62, 2.5, 0.3, 14, True, TEAL); txt(s, 'A V2 não mantém uma quantidade artificial de requisitos apenas para “encher” o documento. A contagem resulta do comportamento do MVP. Termos vagos como “seguro”, “rápido” ou “LGPD compliant” só aparecem quando existe critério ou evidência correspondente.', 3.1, 4.55, 9.1, 0.8, 12, False, DARK); txt(s, 'Referência: docs/ESPECIFICACAO_REQUISITOS_V2.md', 0.75, 6.35, 11.8, 0.3, 10.5, False, MID); footer(s, 4)

# 5
s = blank(); title(s, 'Casos de uso', 'Relações simples, sem setas de fluxo entre ator e caso de uso', 'A modelagem foi separada por ator para reduzir cruzamentos e interpretações ambíguas.')
for x, img, cap in [(0.55,'01_uc_responsavel.png','Responsável'),(4.55,'02_uc_portaria.png','Portaria'),(8.55,'03_uc_gestao.png','Gestão escolar')]:
    rect(s, x, 2.15, 3.65, 4.45, WHITE, BORDER, True); s.shapes.add_picture(str(DIAG / img), Inches(x + 0.16), Inches(2.32), width=Inches(3.33)); txt(s, cap, x + 0.2, 6.13, 3.25, 0.25, 11.5, True, NAVY, PP_ALIGN.CENTER)
footer(s, 5)

# 6
s = blank(); title(s, 'Modelagem', 'Classes, DER e sequência têm responsabilidades diferentes')
for x, img, cap, desc in [(0.45,'04_classes_nucleo.png','Classes conceituais','Conceitos e multiplicidades; FKs não são misturadas ao modelo conceitual.'),(4.55,'06_der_nucleo.png','DER lógico','PK/FK/UQ e cardinalidades; relacionamentos sem setas de processo.'),(8.65,'08_sequencia_presenca.png','Sequência','Setas representam mensagens; retornos são tracejados.')]:
    rect(s, x, 2.12, 3.82, 4.6, WHITE, BORDER, True); s.shapes.add_picture(str(DIAG / img), Inches(x + 0.12), Inches(2.3), width=Inches(3.58)); txt(s, cap, x + 0.18, 5.9, 3.45, 0.26, 12, True, NAVY, PP_ALIGN.CENTER); txt(s, desc, x + 0.2, 6.2, 3.4, 0.42, 9.5, False, MID, PP_ALIGN.CENTER)
footer(s, 6)

# 7
s = blank(); title(s, 'Arquitetura e privacidade', 'O diagrama mostra somente o que realmente existe no MVP')
s.shapes.add_picture(str(DIAG / '09_arquitetura_componentes.png'), Inches(0.65), Inches(2.15), width=Inches(7.0))
bullet_card(s, 8.1, 2.18, 4.55, 1.15, 'Privacidade V2', 'Gestão não lê mensagens privadas de terceiros nem recebe notificações destinadas aos responsáveis.')
bullet_card(s, 8.1, 3.55, 4.55, 1.15, 'Minimização', 'Diretório retorna somente id, nome, perfil e status. Consultas administrativas retornam somente dados necessários.')
bullet_card(s, 8.1, 4.92, 4.55, 1.15, 'Persistência acadêmica', 'Seed sintética versionada; estado mutável em db.runtime.json, criado localmente e ignorado pelo Git.'); footer(s, 7)

# 8
s = blank(); title(s, 'Qualidade', 'Suíte V2: domínio, segurança, API e regressões de privacidade')
metric(s, 0.7, 2.25, 3.2, '24/24', 'testes aprovados localmente — Node.js 22'); metric(s, 4.05, 2.25, 2.65, '0', 'falhas na execução local'); metric(s, 6.85, 2.25, 2.65, '3', 'versões previstas no CI: 18/20/22'); metric(s, 9.65, 2.25, 2.95, '5', 'novas regressões de privacidade')
bullet_card(s, 0.7, 3.75, 3.75, 1.8, 'Regras de domínio', 'Token, RBAC, alunos vinculados, sequência entrada → saída e definição do dia escolar.')
bullet_card(s, 4.8, 3.75, 3.75, 1.8, 'API', 'Login, autorização, presença, CSV, mensagens e feedback.')
bullet_card(s, 8.9, 3.75, 3.7, 1.8, 'Privacidade', 'Diretório mínimo, mensagens isoladas, notificações isoladas, CSV sem DEMO_SEED e reset auditado.')
txt(s, 'Antes do merge final: confirmar o CI da branch e repetir o roteiro funcional.', 0.75, 6.25, 11.7, 0.35, 11, True, WARN, PP_ALIGN.CENTER); footer(s, 8)

# 9
s = blank(); title(s, 'Validação extensionista', 'O projeto não transforma dado histórico em evidência sem comprovação')
bullet_card(s, 0.65, 2.25, 3.8, 2.2, 'O que existe', 'Instrumentos, perfis, protocolo de validação e resultados históricos declarados em versão anterior.')
bullet_card(s, 4.75, 2.25, 3.8, 2.2, 'O que falta', 'Formulários/termos preenchidos, respostas anonimizadas ou CSV real que permita reproduzir quantidade, percentuais e média.', WARN)
bullet_card(s, 8.85, 2.25, 3.8, 2.2, 'Como fechar corretamente', 'Executar tarefas reais, registrar Feedback sem nome do participante, exportar apenas APRESENTACAO e anexar a evidência.')
rect(s, 0.65, 4.85, 12.0, 1.2, LIGHT, BORDER, True); txt(s, 'Status acadêmico correto: PESQUISA DE CAMPO = PARCIAL até existir evidência primária real.', 0.95, 5.17, 11.4, 0.34, 16, True, WARN, PP_ALIGN.CENTER); footer(s, 9)

# 10
s = blank(); title(s, 'Fechamento', 'Como defender o Safe Student depois da revisão V2')
steps = [('1','Abra com o problema e limite do MVP.'),('2','Mostre as 5 funcionalidades essenciais.'),('3','Explique que requisito, classes, DER e sequência são artefatos diferentes.'),('4','Demonstre entrada → notificação → relatório → auditoria.'),('5','Mostre teste inválido e autorização no backend.'),('6','Finalize com a pendência real de evidência de campo e evolução futura.')]
for i, (n, body) in enumerate(steps):
    col = i % 2; row = i // 2; x = 0.75 + col * 6.15; y = 2.2 + row * 1.25
    rect(s, x, y, 5.75, 0.95, WHITE, BORDER, True); txt(s, n, x + 0.15, y + 0.2, 0.5, 0.4, 18, True, TEAL, PP_ALIGN.CENTER); txt(s, body, x + 0.75, y + 0.16, 4.75, 0.52, 11.5, False, DARK)
rect(s, 0.75, 6.15, 11.9, 0.68, NAVY, NAVY, True); txt(s, 'V2.0: mais importante que “parecer completo” é conseguir provar cada afirmação diante da banca.', 1.0, 6.33, 11.4, 0.3, 14, True, WHITE, PP_ALIGN.CENTER); footer(s, 10)

prs.save(OUT / '08_Apresentacao_Final_Safe_Student.pptx')
print('Apresentação V2 gerada.')
